from nltk.tokenize import word_tokenize , sent_tokenize , PunktSentenceTokenizer
from  PyPDF2 import PdfFileReader, PdfFileWriter
from nltk.stem import PorterStemmer
import nltk.data
from nltk.corpus import cgisample, stopwords
import pptx
import docx

#--------------------------------------
#    Variable declaration
#--------------------------------------
ps = PorterStemmer()
stop_words = set(stopwords.words("english"))
MaxMatchWord = []
JunkChar = ['(',')','#','?','.','/']
KeyWords = []
#---------------------------------------------------------
#    This class converts Pdf to Text
#---------------------------------------------------------
class ConvPdfToText():
    def __init__(self,path,path_out,fl_name):
        self.path = path
        self.path_out = path_out
        self.fl_name = fl_name

    def Pdfconv(self):
        newFile = open(self.path_out + self.fl_name + ".txt",mode ='w')
        file = open(self.path,'rb')
        pdfreader = PdfFileReader(file)

        Pagenum = pdfreader.getNumPages()

        for i in range(0,Pagenum):
            Pageobj = pdfreader.getPage(i)
            txt = Pageobj.extractText()
            PageNumber = 'Page Number : '+ str(i) + '\n'
            newFile.writelines(PageNumber)
            newFile.write(txt)
        newFile.close()
        file.close()

#------------------------------------------------
#    Class with Method used to Stem the word Tokens
#-------------------------------------------------

class InputStmmr():
    def __init__(self,input):
        self.input = input

    def SrchpattrnStmmed(self):
        KeyWords =[]
        SrchpattrnTkn = word_tokenize(self.input)
        for token in SrchpattrnTkn:
            if token not in stop_words:
                KeyWords.append(ps.stem(token))
                continue
        #print(KeyWords)
        return KeyWords

#--------------------------------------------------
#     This is the Pattern Matching class
#--------------------------------------------------
class TextMinning(InputStmmr,ConvPdfToText):
    def __init__(self, path_out ,Text_comp,fl_name,file ):
        self.path_out = path_out
        self.Text_comp = Text_comp
        self.fl_name = fl_name
        self.file = file
    def Textminner(self):
        #print('Textminner fl_name: '+ self.fl_name)
        if self.file.endswith(".pdf"):
            Txtfile = open(self.path_out + self.fl_name + '.txt' , mode="r")
            lines = Txtfile.readlines()
            SrchWrd =[]
            Output_Match = []

            for line in lines:
                Objline = InputStmmr(line)
                wcnt = 0
                lineSrch = Objline.SrchpattrnStmmed()
                for w in self.Text_comp:

                    if w not in JunkChar and w in lineSrch :
                        wcnt+=1
                        MaxMatchWord.append(w)
                        SrchWrd.append(w)
                        if len(set(self.Text_comp))==1 or (len(set(self.Text_comp)) - wcnt)<=1:
                            Output_Match.append(line)
                    #    continue
        elif self.file.endswith(".docx") or self.file.endswith(".doc") or self.file.endswith(".pptx"):
            with open(self.path, 'rb') as f:
                doc = docx.Document(f)
            fullText = ''
            for para in doc.paragraphs:
                wcnt = 0
                Objline = InputStmmr(para)
                lineSrch = Objline.SrchpattrnStmmed()

                for w in self.Text_comp:
                    if w not in JunkChar and w in lineSrch:
                        wcnt += 1
                        MaxMatchWord.append(w)
                        SrchWrd.append(w)
                    Output_Match.append(line)

        if len(set(Output_Match))!= 0:
            print("Information Found In File: ")
            print(self.path_out+self.fl_name+'.pdf')
            for i in range(len(Output_Match)):

                print(Output_Match[i])

#------------------------------------------------
#       This class converts Docx and Pptx to Text
#------------------------------------------------
class ReadDocxPptx(TextMinning):
    def __init__(self, path,path_out, fl_name,Srchpattrn ):
        self.path = path
        self.path_out = path_out
        self.fl_name = fl_name
        self.Srchpattrn = Srchpattrn
    def ConvDocToTxt(self):
        #newFile = open(self.path_out + self.fl_name + ".txt", mode='w')
        print('fl_name: ' + self.fl_name)
        # with open(self.path ,'rb') as f:
        #     doc = docx.Document(f)
        # fullText = ''
        # for para in doc.paragraphs:
        #     fullText = para
        #     ObjInputStmmr = InputStmmr(self.Srchpattrn)
        #     ObjStemmr = TextMinning(self.path_out, fullText, self.fl_name)
        #     ObjInputStmmr.SrchpattrnStmmed()
        #     ObjStemmr.Textminner()
        print("I am in word doc read .......")
        # f.close()

    def ConvPptxToTxt(self):
        txt = ' '
        prs = pptx.Presentation(self.path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    txt = txt + shape.text_frame.text
                    print("I am in PPTX......")
        #newFile.close()
        return txt
